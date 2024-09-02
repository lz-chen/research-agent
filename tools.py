from pptx import Presentation
from pptx.util import Inches


def get_all_layouts_info(presentation_file):
    """
    Extracts detailed information from all slide layouts within a PowerPoint presentation template file,
    including layout name, placeholders, number of shapes, background presence, slide master name,
    and font sizes for text placeholders.

    Args:
        presentation_file (str): The path to the PowerPoint presentation file (e.g., 'template.pptx').

    Returns:
        list of dict: A list of dictionaries, each containing detailed information about a slide layout:
            - 'name' (str): The name of the slide layout.
            - 'placeholders' (list of dict): A list of dictionaries, each containing:
                - 'index' (int): The index of the placeholder within the layout.
                - 'name' (str): The name of the placeholder.
                - 'shape_type' (str): The type of shape (e.g., TEXT_BOX) for the placeholder.
                - 'font_size' (float or None): The font size in points (pt) for the first paragraph in the placeholder,
                  or None if the font size is not explicitly defined.
            - 'number_of_shapes' (int): The total number of shapes in the layout.
            - 'has_background' (bool): Whether the layout has a background defined.
            - 'slide_master_name' (str): The name of the slide master associated with this layout.

    Example:
        layouts_info_list = get_all_layouts_info('template.pptx')
        print(layouts_info_list)

        Output might look like:
        [
            {
                'name': 'Title Slide',
                'placeholders': [
                    {'index': 0, 'name': 'Title 1', 'shape_type': 'TEXT_BOX', 'font_size': 44.0},
                    {'index': 1, 'name': 'Subtitle 2', 'shape_type': 'TEXT_BOX', 'font_size': 32.0}
                ],
                'number_of_shapes': 3,
                'has_background': True,
                'slide_master_name': 'Office Theme'
            },
            {
                'name': 'Title and Content',
                'placeholders': [
                    {'index': 0, 'name': 'Title 1', 'shape_type': 'TEXT_BOX', 'font_size': 32.0},
                    {'index': 1, 'name': 'Content Placeholder 2', 'shape_type': 'TEXT_BOX', 'font_size': None}
                ],
                'number_of_shapes': 5,
                'has_background': False,
                'slide_master_name': 'Office Theme'
            },
            # More layout dictionaries...
        ]
    """
    # Load the presentation
    prs = Presentation(presentation_file)

    # List to store information about all slide layouts
    layouts_info = []

    # Iterate over each layout in the presentation
    for layout in prs.slide_layouts:
        # Create a dictionary to store the layout information
        layout_info = {}

        # Get the name of the layout
        layout_info['layout_name'] = layout.name

        # Get information about placeholders
        placeholders_info = []
        for placeholder in layout.placeholders:
            text_frame = placeholder.text_frame
            font_size = None
            if text_frame and text_frame.paragraphs and text_frame.paragraphs[0].font.size:
                font_size = text_frame.paragraphs[0].font.size.pt  # Convert to points (pt)
            auto_size = text_frame.auto_size.name if text_frame.auto_size else None
            placeholder_info = {
                'index': placeholder.placeholder_format.idx,
                'name': placeholder.name,
                'shape_type': placeholder.shape_type,
                'left': Inches(placeholder.left.inches),
                'top': Inches(placeholder.top.inches),
                'width': Inches(placeholder.width.inches),
                'height': Inches(placeholder.height.inches),
                'font_size': font_size,
                'auto_size': auto_size
            }
            placeholders_info.append(placeholder_info)
        layout_info['placeholders'] = placeholders_info

        # Get the number of shapes in the layout
        layout_info['number_of_shapes'] = len(layout.shapes)

        # Check if the layout has a background
        layout_info['has_background'] = layout.background is not None

        # Get the name of the slide master
        layout_info['slide_master_name'] = layout.slide_master.name

        # Add the layout information to the list
        layouts_info.append(layout_info)

    return layouts_info

# info = get_all_layouts_info('/Users/lzchen/PycharmProjects/research_agent/data/Inmeta Brand guidelines 2023.pptx')
# from pprint import pprint
# pprint(info)

# def outlines2code(slide_outlines: str):
#     llm = llms.new_gpt4o(temperature=0.0)
